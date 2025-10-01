from fastapi import FastAPI, UploadFile, File, Form, Query, HTTPException
from fastapi.responses import JSONResponse
from typing import List, Optional
import os
from dotenv import load_dotenv
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
from azure.ai.formrecognizer import DocumentAnalysisClient
from openai import AzureOpenAI
import pandas as pd
from docx import Document
from pptx import Presentation
import PyPDF2
from datetime import datetime, timedelta, UTC
import requests
from fastapi.middleware.cors import CORSMiddleware
from io import BytesIO
import re

load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://red-desert-05e26cb03.1.azurestaticapps.net"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Azure setup (same as before).
blob_connection_string = "DefaultEndpointsProtocol=https;AccountName=contracts;AccountKey=0KBsuwH1J8/ei4/su4gmh0kOTVqXfrCGhgD6C/cpM3WpuvXEK8Botw6OItcmy5cZ0048ATewyzPK+ASt64xYvg==;EndpointSuffix=core.windows.net"
container_name = "data"
blob_service_client = BlobServiceClient.from_connection_string(blob_connection_string)
container_client = blob_service_client.get_container_client(container_name)

search_endpoint = "https://contract-analyser.search.windows.net"
search_key = os.getenv("CONTRACT_SEARCH_KEY")
search_index = "azureblob-index"
search_client = SearchClient(
    endpoint=search_endpoint,
    index_name=search_index,
    credential=AzureKeyCredential(search_key)
)

endpoint = "https://pilot-contract.openai.azure.com/"
deployment = "gpt-4o"
subscription_key = os.getenv("CONTRACT_ANALYSIS_APIKYEY")
client = AzureOpenAI(
    azure_endpoint=endpoint,
    api_key=subscription_key,
    api_version="2024-12-01-preview",
)

search_service = "contract-analyser"
#search_admin_key = os.getenv("CONTRACT_SEARCH_KEY")
indexer_name = "azureblob-indexer"

dockey = os.getenv("CONTRACT_DOCINTELL_KEY")
docendpoint = "https://contract2.cognitiveservices.azure.com/"


# Utility functions (same as before, but no Streamlit calls)
def analyze_read(url):
    document_analysis_client = DocumentAnalysisClient(
        endpoint=docendpoint, credential=AzureKeyCredential(dockey)
    )
    poller = document_analysis_client.begin_analyze_document_from_url(
        "prebuilt-layout", url)
    result = poller.result()
    # Extract tables as text
    tables = []
    for table in result.tables:
        rows = []
        for row_idx in range(table.row_count):
            row = []
            for col_idx in range(table.column_count):
                cell = next(
                    (c for c in table.cells if c.row_index == row_idx and c.column_index == col_idx), None
                )
                row.append(cell.content if cell else "")
            rows.append(" | ".join(row))
        tables.append("\n".join(rows))
    return "\n\n".join(tables) if tables else result.content

def run_indexer_once():
    endpoint = f"https://{search_service}.search.windows.net/indexers/{indexer_name}/run?api-version=2023-11-01"
    headers = {"api-key": search_key, "Content-Type": "application/json"}
    response = requests.post(endpoint, headers=headers)
    if response.status_code == 202:
        return {"status": "success", "message": "Indexer run triggered."}
    else:
        return {"status": "error", "message": response.text}

def retrieve_context_from_vector_store(query, top_k=5):
    results = search_client.search(query, top=top_k)
    context_chunks = []
    for result in results:
        # Adjust 'content' to match your index field name
        context_chunks.append(result['content'])
    return "\n\n".join(context_chunks)

def chunk_text(text, chunk_size):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def extract_txt(file):
    return file.read().decode("utf-8")

def extract_xlsx(file):
    # Ensure file is a BytesIO object for pandas
    if not hasattr(file, 'read'):
        file = BytesIO(file)
    xls = pd.ExcelFile(file)
    all_text = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        rows = df.astype(str).apply(lambda row: ' | '.join(row), axis=1).tolist()
        all_text.append(f"--- Sheet: {sheet_name} ---")
        all_text.extend(rows)
    return '\n'.join(all_text)

# New: Extract CSV
def extract_csv(file_bytes):
    # file_bytes is bytes, so wrap in BytesIO
    df = pd.read_csv(BytesIO(file_bytes))
    rows = df.astype(str).apply(lambda row: ' | '.join(row), axis=1).tolist()
    return '\n'.join(rows)

def extract_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_docx(file_bytes):
    doc = Document(BytesIO(file_bytes))
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_pptx(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
            # Extract tables from shapes
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    text.append(" | ".join(row_text))
            # Extract tables from grouped shapes
            if hasattr(shape, "shapes"):
                for subshape in shape.shapes:
                    if hasattr(subshape, "has_table") and subshape.has_table:
                        table = subshape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            text.append(" | ".join(row_text))
    return "\n".join(text)

def get_blob_sas_url(blob_client):
    sas_token = generate_blob_sas(
        account_name=blob_client.account_name,
        container_name=blob_client.container_name,
        blob_name=blob_client.blob_name,
        account_key=blob_service_client.credential.account_key,
        permission=BlobSasPermissions(read=True),
        expiry=datetime.now(UTC) + timedelta(hours=1)  # Use timezone-aware UTC
    )
    return f"{blob_client.url}?{sas_token}"

# FastAPI endpoints

@app.get("/documents")
def list_documents():
    blob_list = container_client.list_blobs()
    doc_set = set()
    doc_map = {}
    for blob in blob_list:
        name = blob.name
        if "/" in name:
            classification, doc_name = name.split("/", 1)
        else:
            classification, doc_name = "unknown", name
        # Remove _partX.txt if present
        base_name = doc_name.split("_part")[0] if "_part" in doc_name else doc_name
        doc_set.add(base_name)
        doc_map[base_name] = classification
    doc_data = [{"Document": doc, "Classification": doc_map[doc]} for doc in doc_set]
    return {"documents": doc_data}

@app.delete("/documents")
def delete_documents(doc_names: List[str] = Query(...)):
    blob_list = list(container_client.list_blobs())
    deleted = []
    for base_name in doc_names:
        for blob in blob_list:
            if "/" in blob.name:
                _, doc_name = blob.name.split("/", 1)
            else:
                doc_name = blob.name
            if doc_name.startswith(base_name):
                container_client.delete_blob(blob.name)
                deleted.append(blob.name)
    run_indexer_once()
    return {"deleted": deleted}

@app.post("/upload")
async def upload_documents(files: List[UploadFile] = File(...), categories: List[str] = Form(...)):
    chunk_size = 4000
    uploaded = []
    for uploaded_file, category in zip(files, categories):
        print("Received file:", uploaded_file.filename)
        ext = uploaded_file.filename.split('.')[-1].lower()
        file_bytes = await uploaded_file.read()
        if ext == "xlsx":
            try:
                text = extract_xlsx(file_bytes)
            except Exception as e:
                return JSONResponse(status_code=400, content={"error": f"Failed to process XLSX: {str(e)}"})
        elif ext == "csv":
            text = extract_csv(file_bytes)
        elif ext == "pptx":
            text = extract_pptx(file_bytes)
        elif ext == "docx":
            text = extract_docx(file_bytes)
        else:
            blob_name = f"{category}/{uploaded_file.filename}"
            blob_client = container_client.get_blob_client(blob_name)
            blob_client.upload_blob(file_bytes, overwrite=True, metadata={"classification": category})
            blob_url = get_blob_sas_url(blob_client)
            try:
                text = analyze_read(blob_url)
            except Exception as e:
                container_client.delete_blob(blob_name)
                continue
            container_client.delete_blob(blob_name)
        parts = chunk_text(text, chunk_size)
        for idx, part in enumerate(parts):
            chunk_blob_name = f"{category}/{uploaded_file.filename}_part{idx+1}.txt"
            chunk_blob_client = container_client.get_blob_client(chunk_blob_name)
            chunk_blob_client.upload_blob(part, overwrite=True, metadata={"classification": category})
        uploaded.append(uploaded_file.filename)
    run_indexer_once()
    return {"uploaded": uploaded}

@app.post("/chat")
def chat(messages: List[dict]):
    # Get the latest user question
    user_message = next((m for m in reversed(messages) if m["role"] == "user"), None)
    user_question = user_message["content"] if user_message else ""
    
    # Retrieve context ONLY from stored documents (no web search or external sources)
    context = retrieve_context_from_vector_store(user_question, top_k=5)

    # System prompt strictly restricts answers to stored context
    rag_messages = [
        {
            "role": "system",
            "content": (
                """
                You are a lead document reviewer. Only answer using the provided document context below. Do NOT use any external sources, web search, or your own general knowledge. If the answer is not found in the context, reply with: 'I sorry I don't have any context with that detail.'
                Always provide the reference to the source document or section and provide a link if possible. If the user asks for comparisons between classifications, only use the context provided for each classification..
                """
            ),
        },
        {
            "role": "user",
            "content": f"Context:\n{context}\n\nQuestion: {user_question}"
        }
    ]

    completion = client.chat.completions.create(
        model=deployment,
        messages=rag_messages,
        max_tokens=800,
        temperature=0.2,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None,
        stream=False
    )
    response = completion.choices[0].message.content

    # Remove HTML tags only (for safety)
    response = re.sub(r'<[^>]+>', '', response)

    return {"response": response}

@app.post("/run_indexer")
def run_indexer():
    result = run_indexer_once()
    return result
